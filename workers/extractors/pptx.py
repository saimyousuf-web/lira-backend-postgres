"""
PPTX chunking + embedding extractor — Postgres version.

Migrated from DynamoDB (chunks_table.put_item + file_table.update_item ×2)
to async SQLAlchemy, following the exact same pattern as the PDF/DOCX/Storyline
extractors. Only the persistence layer changed; all PPTX-specific logic
(image extraction, Textract OCR, text chunking) is untouched.
"""

import io
import json
import os
import re
import uuid
from datetime import datetime, timezone

import boto3
from fastapi import HTTPException
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
from qdrant_client.models import PointStruct
from sqlalchemy import select
from sqlalchemy.exc import IntegrityError, SQLAlchemyError
from sqlalchemy.ext.asyncio import AsyncSession

from models.chunks import Chunk
from models.course_module import CourseModule
from models.module import Module
from utils.text_chunking import chunk_text_with_images, group_images_with_previous_text
from workers.vector import generate_embedding, normalize_vector
from workers.qdrnt_vector import get_qdrant_client, QDRANT_COLLECTION

# ---------------------------------------------------------------------------
# Image quality filters
# ---------------------------------------------------------------------------
MIN_WIDTH          = 50
MIN_HEIGHT         = 50
MIN_AREA           = 13_000
MIN_BYTES          = 18_000
MAX_BYTES          = 2_500_000
MIN_BYTES_PER_PIXEL = 0.15

EMU_PER_INCH = 914_400
DPI          = 96

REGION = os.environ.get("REGION")
s3       = boto3.client("s3",       region_name=REGION)
textract = boto3.client("textract", region_name=REGION)


# ===========================================================================
# Shape iterator
# ===========================================================================

def iter_shapes(shapes):
    """Recursively yield shapes, flattening groups."""
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes)


# ===========================================================================
# S3 + Textract helpers
# ===========================================================================

def _upload_slide_images(prs: Presentation, s3_prefix: str, S3_BUCKET: str) -> int:
    """Extract qualifying images from every slide and upload to S3."""
    saved = skipped = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_saved = 0

        for shape in iter_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            img_bytes = shape.image.blob
            ext       = shape.image.ext.lower()
            size      = len(img_bytes)
            w         = int(shape.width  * DPI / EMU_PER_INCH)
            h         = int(shape.height * DPI / EMU_PER_INCH)
            area      = w * h

            if (
                w    < MIN_WIDTH
                or h < MIN_HEIGHT
                or area < MIN_AREA
                or size < MIN_BYTES
                or size > MAX_BYTES
                or (size / max(area, 1)) < MIN_BYTES_PER_PIXEL
            ):
                skipped += 1
                continue

            slide_saved += 1
            saved       += 1
            filename     = f"slide{slide_index}_img{slide_saved}.{ext}"

            s3.put_object(
                Bucket=S3_BUCKET,
                Key=s3_prefix + filename,
                Body=img_bytes,
                ContentType=f"image/{ext}",
            )

    print(f"Uploaded {saved} PPTX images to S3 (skipped {skipped})")
    return saved


def _load_slide_images_from_s3(s3_prefix: str, S3_BUCKET: str) -> dict[int, list[dict]]:
    """
    List objects under s3_prefix, run Textract OCR on each.
    Returns {slide_number: [{filename, text, s3_key}, ...]}.
    """
    slide_images: dict[int, list[dict]] = {}
    pattern = re.compile(r"slide(\d+)_", re.IGNORECASE)

    paginator = s3.get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=s3_prefix):
        for obj in page.get("Contents") or []:
            key      = obj["Key"]
            filename = key.split("/")[-1]
            match    = pattern.match(filename)
            if not match:
                continue

            slide_num    = int(match.group(1))
            ocr_text     = ""
            try:
                resp = textract.detect_document_text(
                    Document={"S3Object": {"Bucket": S3_BUCKET, "Name": key}}
                )
                ocr_text = "\n".join(
                    b.get("Text", "")
                    for b in resp.get("Blocks", [])
                    if b.get("BlockType") == "LINE" and b.get("Text")
                )
            except Exception as e:
                print(f"Textract OCR failed for {key}: {e}")

            slide_images.setdefault(slide_num, []).append(
                {"filename": filename, "text": ocr_text, "s3_key": key}
            )

    for imgs in slide_images.values():
        imgs.sort(key=lambda x: x["filename"])

    return slide_images


def _build_all_slide_chunks(
    prs: Presentation,
    slide_images: dict[int, list[dict]],
) -> list[dict]:
    """
    For each slide: extract digital text, merge with OCR'd image text,
    split into overlapping chunks.

    Returns a flat list:
        [{chunk_id, text, image (list[str] filenames), page_num}, ...]
    """
    flat_chunks: list[dict] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_num = slide_index

        # --- Digital text ---
        text_blocks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_blocks.append(t)

        # --- Build content list ---
        content: list[dict] = [{"text": t, "image": []} for t in text_blocks]
        for img in slide_images.get(slide_num, []):
            content.append({
                "text":        "",
                "image":       [img],
                "image_texts": [img.get("text", "")],
            })

        content = group_images_with_previous_text(content)

        # Flatten image dicts → filename strings
        for item in content:
            item["image"] = [
                img["filename"] if isinstance(img, dict) else img
                for img in item.get("image", [])
            ]

        # Sliding-window split
        for item in content:
            for chunk in chunk_text_with_images(
                item["text"], item["image"], max_chars=1200, overlap=150
            ):
                flat_chunks.append({**chunk, "page_num": slide_num})

    return flat_chunks


# ===========================================================================
# Qdrant embed helper
# ===========================================================================

def _embed_and_upsert_qdrant(
    chunk_id:    uuid.UUID,
    text:        str,
    org_id:      str,
    course_id:   uuid.UUID,
    course_name: str,
    module_id:   uuid.UUID,
    page_num:    int,
) -> tuple[str, int]:
    try:
        embedding = generate_embedding(text)

        if not (isinstance(embedding, list) and len(embedding) == 1024):
            return "failed", 0

        normalized = normalize_vector(embedding)
        get_qdrant_client().upsert(
            collection_name=QDRANT_COLLECTION,
            points=[
                PointStruct(
                    id=str(chunk_id),
                    vector=normalized,
                    payload={
                        "chunk_id":        str(chunk_id),
                        "organization_id": org_id,
                        "course_id":       str(course_id),
                        "course_name":     course_name,
                        "module_id":       str(module_id),
                        "slide_index":     page_num,
                        "chunk_index":     1,
                    },
                )
            ],
        )
        return "completed", len(embedding)

    except Exception as e:
        print(f"Embedding failed for chunk {chunk_id}: {e}")
        return "failed", 0


# ===========================================================================
# Main async extractor
# ===========================================================================

async def chunk_and_embed_pptx(
    pptx_bytes:  bytes,
    ctx_orgid:   str,
    ctx_cid:     str,
    course_name: str,
    ctx_moid:    str,
    S3_BUCKET:   str,
    db:          AsyncSession,
):
    """
    Extract, chunk, embed, and persist a PPTX module.

    Pipeline
    --------
    1. Validate Module + CourseModule exist in Postgres.
    2. Upload qualifying images → S3, run Textract OCR.
    3. Build overlapping text + image chunks per slide.
    4. Embed each chunk (bge-large-en-v1.5) → upsert Qdrant.
    5. Persist Chunk rows + update Module / CourseModule — one transaction.
    """
    user_id = uuid.UUID("6418e458-50a1-70fe-9d3e-b52f5d2df57c")

    try:
        course_id = uuid.UUID(ctx_cid)
        module_id = uuid.UUID(ctx_moid)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid UUID format")

    # ------------------------------------------------------------------
    # 1. Validate Module + CourseModule
    # ------------------------------------------------------------------
    try:
        mod_result = await db.execute(select(Module).where(Module.id == module_id))
        module: Module | None = mod_result.scalar_one_or_none()
        if not module:
            raise HTTPException(status_code=404, detail="Module not found")

        cm_result = await db.execute(
            select(CourseModule).where(
                CourseModule.cid == course_id,
                CourseModule.moid == module_id,
            )
        )
        course_module: CourseModule | None = cm_result.scalar_one_or_none()
        if not course_module:
            raise HTTPException(status_code=404, detail="CourseModule association not found")

    except HTTPException:
        raise
    except SQLAlchemyError as e:
        raise HTTPException(status_code=500, detail=f"Database error: {e}")

    # ------------------------------------------------------------------
    # 2. Extract images → S3 → Textract OCR
    # ------------------------------------------------------------------
    s3_prefix = f"material/{ctx_orgid}/pptx/{ctx_cid}/{ctx_moid}/images/"
    prs       = Presentation(io.BytesIO(pptx_bytes))

    try:
        _upload_slide_images(prs, s3_prefix, S3_BUCKET)
        slide_images = _load_slide_images_from_s3(s3_prefix, S3_BUCKET)
    except Exception as e:
        print(f"Image pipeline failed (continuing without images): {e}")
        slide_images = {}

    # ------------------------------------------------------------------
    # 3. Build chunks
    # ------------------------------------------------------------------
    flat_chunks = _build_all_slide_chunks(prs, slide_images)

    if not flat_chunks:
        raise HTTPException(
            status_code=422, detail="No content could be extracted from PPTX"
        )

    # ------------------------------------------------------------------
    # 4 + 5. Embed → Qdrant, persist Chunk rows
    # ------------------------------------------------------------------
    failed_count = 0

    for chunk_data in flat_chunks:
        chunk_id: uuid.UUID = chunk_data["chunk_id"]
        text:     str       = chunk_data["text"]
        img_keys: list[str] = chunk_data["image"]
        page_num: int       = chunk_data["page_num"]

        status, _ = _embed_and_upsert_qdrant(
            chunk_id=chunk_id,
            text=text,
            org_id=ctx_orgid,
            course_id=course_id,
            course_name=course_name,
            module_id=module_id,
            page_num=page_num,
        )

        if status == "failed":
            failed_count += 1
            continue

        db.add(Chunk(
            id=chunk_id,
            moid=module_id,
            cid=course_id,
            txt=text,
            # img_keys=img_keys or None,   # uncomment after adding column to Chunk
            crtby=user_id,
            updby=user_id,
        ))

    # ------------------------------------------------------------------
    # Mark Module + CourseModule processed
    # ------------------------------------------------------------------
    now = datetime.now(timezone.utc)

    module.isvec = True
    module.sts   = "APPROVED"
    module.updat = now
    module.updby = user_id

    course_module.updat = now
    course_module.updby = user_id

    # ------------------------------------------------------------------
    # Single commit
    # ------------------------------------------------------------------
    try:
        await db.commit()
    except IntegrityError as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Integrity error: {e.orig}")
    except SQLAlchemyError as e:
        await db.rollback()
        try:
            module.sts   = "FAILED"
            module.updat = now
            await db.commit()
        except Exception:
            pass
        raise HTTPException(status_code=500, detail=f"Database error: {e}")

    total = len(flat_chunks)
    saved = total - failed_count
    print(f"PPTX chunking + embedding complete. {saved}/{total} chunks saved.")

    return {
        "status":        "success",
        "status_code":   200,
        "message":       "PPTX chunked and embedded successfully",
        "chunks_total":  total,
        "chunks_saved":  saved,
        "chunks_failed": failed_count,
    }